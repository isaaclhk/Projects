# import necessary libraries
from typing import List

from docx import Document as doc_reader
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_loaders.csv_loader import CSVLoader
from langchain_community.vectorstores import Chroma
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from pptx import Presentation

# Load and extract text from PDF using PyPDFLoader
def extract_text_from_pdf(file_path: str) -> List[Document]:
    '''
    Load and extract text from a PDF file.
    Args:
        file_path (str): Path to the PDF file.
    Returns:
        List[Document]: A list of Document objects containing the text from the PDF.
    '''
    loader = PyPDFLoader(file_path)
    documents = loader.load()
    return documents

# Split the text into chunks
def split_documents(
    documents, chunk_size: int=500, chunk_overlap:int=0
    ) -> List[Document]:
    '''
    Split documents into smaller chunks for processing.
    
    Args:
        documents (List[Document] or str): The documents to split.
        chunk_size (int): The size of each chunk.
        chunk_overlap (int): The overlap between chunks.
        
    Returns:
        List[Document]: A list of Document objects containing the split text.
    '''
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )
    
    # Wrap raw text in Document objects
    if isinstance(documents, str):
        documents = [Document(page_content=documents)]
    return splitter.split_documents(documents)

def load_docx(file_path: str) -> str:
    '''
    Load and extract text from a DOCX file.
    
    Args:
        file_path (str): Path to the DOCX file.
    
    Returns:
        str: Extracted text from the DOCX file.
    '''
    doc = doc_reader(file_path)
    documents = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    return documents

def load_csv(file_path: str) -> List[Document]:
    '''
    Load and extract text from a CSV file.
    Args:
        file_path (str): Path to the CSV file.
    Returns:
        List[Document]: A list of Document objects containing the text from the CSV.
    '''
    loader = CSVLoader(file_path)
    documents = loader.load()
    return documents

def extract_text_from_pptx(file_path):
    '''
    Load and extract text from a PowerPoint (PPTX) file.
    Args:
        file_path (str): Path to the PPTX file.
    Returns:
        str: Extracted text from the PPTX file.
    '''
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Embed the chunks using MedEmbed
def embed_chunks(chunks: List[Document], embed_model: str) -> Chroma:
    '''
    Embed the text chunks using the MedEmbed model.
    Args:
        chunks (List[Document]): List of Document objects containing text chunks.
        embed_model (str): The model to use for embedding.
    Returns:
        Chroma: A vector store containing the embedded chunks.
    '''
    embedding_model = HuggingFaceEmbeddings(model_name=embed_model)
    vectorstore = Chroma.from_documents(chunks, embedding_model)
    return vectorstore

# Perform similarity search
def search_similar_chunks(
    query: str, vectorstore: Chroma, k: int=3
    ) -> List[Document]:
    return vectorstore.similarity_search(query, k=k)